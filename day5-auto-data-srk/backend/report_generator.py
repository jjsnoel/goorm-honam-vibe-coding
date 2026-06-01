"""Markdown and PPT report generation."""

from __future__ import annotations

import io
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt


def build_markdown_report(
    eda: dict[str, Any],
    analysis: dict[str, Any],
) -> str:
    s = eda["structure"]
    m = eda["missing"]
    o = eda["outliers"]
    c = eda["correlation"]
    now = datetime.now().strftime("%Y-%m-%d %H:%M")

    lines = [
        f"# 데이터 분석 보고서: {eda['filename']}",
        "",
        f"> 생성일시: {now} | Data Explorer AI",
        "",
        "## Executive Summary",
        "",
        f"- **데이터 규모**: {s['rows']:,}행 × {s['columns']}열",
        f"- **결측률**: {m['overall_missing_rate']}%",
        f"- **수치형 변수**: {eda['distributions']['numeric_count']}개",
        f"- **AI 분석**: {'GPT 보강' if analysis.get('ai_enhanced') else '규칙 기반'}",
        "",
        "---",
        "",
        "## 1. 데이터 구조 분석",
        "",
        f"| 항목 | 값 |",
        f"|------|-----|",
        f"| 행 수 | {s['rows']:,} |",
        f"| 열 수 | {s['columns']} |",
        f"| 메모리 | {s['memory_mb']} MB |",
        f"| 중복 행 | {s['duplicate_rows']:,} |",
        "",
        "### 컬럼 정보",
        "",
        "| 컬럼 | 타입 | 고유값 | 비결측 |",
        "|------|------|--------|--------|",
    ]

    for col in s["column_details"][:30]:
        lines.append(
            f"| {col['name']} | {col['dtype']} | {col['unique']} | {col['non_null']} |"
        )

    lines.extend(
        [
            "",
            "---",
            "",
            "## 2. 결측치 분석",
            "",
            f"- 전체 결측 셀: {m['total_missing_cells']:,} / {m['total_cells']:,}",
            f"- 결측 컬럼 수: {m['columns_with_missing']}",
            f"- 완전 행 수: {m['complete_rows']:,}",
            "",
            "| 컬럼 | 결측 수 | 결측률(%) |",
            "|------|---------|-----------|",
        ]
    )

    for col in m["by_column"][:20]:
        if col["missing_count"] > 0:
            lines.append(
                f"| {col['column']} | {col['missing_count']} | {col['missing_rate']} |"
            )

    lines.extend(
        [
            "",
            "---",
            "",
            "## 3. 이상치 분석 (IQR / Z-score)",
            "",
            f"- 분석 수치형 변수: {o['numeric_column_count']}개",
            f"- IQR 이상치 총합: {o['total_iqr_outliers']:,}",
            "",
            "| 컬럼 | IQR 이상치 | IQR 비율(%) | Z-score 이상치 |",
            "|------|------------|-------------|----------------|",
        ]
    )

    for item in o["columns_analyzed"][:15]:
        lines.append(
            f"| {item['column']} | {item['iqr']['count']} | {item['iqr']['rate']} | "
            f"{item['zscore']['count']} |"
        )

    lines.extend(["", "---", "", "## 4. 변수 분포 분석", ""])

    for stat in eda["distributions"]["numeric_summary"][:10]:
        lines.append(
            f"**{stat['column']}**: mean={stat.get('mean')}, std={stat.get('std')}, "
            f"skew={stat.get('skewness')}"
        )

    for cat in eda["distributions"]["categorical_summary"][:5]:
        top = cat["top_values"][0] if cat["top_values"] else {}
        lines.append(
            f"**{cat['column']}** (unique={cat['unique']}): "
            f"최빈값={top.get('value', 'N/A')} ({top.get('count', 0)})"
        )

    lines.extend(["", "---", "", "## 5. 상관관계 분석", ""])

    if c.get("available"):
        for pair in c.get("top_pairs", [])[:10]:
            lines.append(
                f"- {pair['x']} ↔ {pair['y']}: **{pair['correlation']:.4f}**"
            )
    else:
        lines.append(c.get("message", "상관관계 분석 불가"))

    lines.extend(["", "---", "", "## 6. 주요 인사이트 10개", ""])
    for i, insight in enumerate(analysis["insights"], 1):
        lines.append(f"{i}. {insight}")

    ml = analysis["ml_strategy"]
    lines.extend(["", "---", "", "## 7. 머신러닝 적용 전략", ""])

    if ml.get("ai_summary"):
        lines.extend([ml["ai_summary"], ""])

    lines.append("### 전처리")
    for p in ml["recommended_preprocessing"]:
        lines.append(f"- {p}")

    lines.append("\n### 추천 모델")
    for model in ml["recommended_models"]:
        lines.append(f"- {model}")

    lines.append(f"\n### 검증 전략\n- {ml['validation_strategy']}")

    lines.append("\n### 다음 단계")
    for step in ml["next_steps"]:
        lines.append(f"- {step}")

    lines.extend(["", "---", "", "## 8. PPT 보고서 초안", ""])
    for slide in build_ppt_outline(eda, analysis):
        lines.append(f"### Slide {slide['number']}: {slide['title']}")
        for bullet in slide["bullets"]:
            lines.append(f"- {bullet}")
        lines.append("")

    return "\n".join(lines)


def build_ppt_outline(
    eda: dict[str, Any],
    analysis: dict[str, Any],
) -> list[dict[str, Any]]:
    s = eda["structure"]
    m = eda["missing"]

    return [
        {
            "number": 1,
            "title": "Executive Summary",
            "bullets": [
                f"데이터셋: {eda['filename']}",
                f"{s['rows']:,}행 × {s['columns']}열",
                f"결측률 {m['overall_missing_rate']}%",
                "자동 EDA 기반 핵심 발견 요약",
            ],
        },
        {
            "number": 2,
            "title": "Dataset Overview",
            "bullets": [
                f"수치형 {eda['distributions']['numeric_count']}개 / "
                f"범주형 {eda['distributions']['categorical_count']}개",
                f"중복 행 {s['duplicate_rows']:,}개",
                f"메모리 {s['memory_mb']} MB",
            ],
        },
        {
            "number": 3,
            "title": "Data Quality Assessment",
            "bullets": [
                f"결측 컬럼 {m['columns_with_missing']}개",
                f"완전 행 {m['complete_rows']:,}개",
                "결측치 히트맵 / 바 차트 참조",
            ],
        },
        {
            "number": 4,
            "title": "Outlier Analysis",
            "bullets": [
                f"IQR 이상치 총 {eda['outliers']['total_iqr_outliers']:,}개",
                "고이상치 변수: "
                + (", ".join(eda["outliers"]["high_outlier_columns"][:5]) or "없음"),
            ],
        },
        {
            "number": 5,
            "title": "Feature Distribution",
            "bullets": [
                "히스토그램 & 박스플롯으로 분포 확인",
                "왜도/첨도 기반 변환 필요성 검토",
            ],
        },
        {
            "number": 6,
            "title": "Correlation Analysis",
            "bullets": [
                "상관 히트맵으로 변수 관계 파악",
                "다중공선성 위험 변수 식별",
            ],
        },
        {
            "number": 7,
            "title": "Key Findings",
            "bullets": analysis["insights"][:5],
        },
        {
            "number": 8,
            "title": "ML Strategy & Recommendations",
            "bullets": analysis["ml_strategy"]["recommended_models"][:3]
            + analysis["ml_strategy"]["next_steps"][:2],
        },
    ]


def build_pptx_bytes(eda: dict[str, Any], analysis: dict[str, Any]) -> bytes:
    prs = Presentation()
    outline = build_ppt_outline(eda, analysis)

    for slide_data in outline:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data["title"]
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(slide_data["bullets"]):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
